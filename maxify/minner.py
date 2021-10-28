from pptx import Presentation
from campaign import ExceptionHander
import os


class TextMinner:
    def __init__(self, file_path):
        self.exception_handler = ExceptionHander()
        self.presentation = None
        if isinstance(file_path, str):
            try:
                self.file = open(file_path)
                self.presentation = Presentation(self.file)
            except Exception as e:
                self.exception_handler.handle_exception(e)
        else:
            raise ValueError("Invalid file for Presentation")

    def extract_text(self):
        ret_value = []
        try:
            # for slide in self.presentation.slides:
            #     for shape in slide.shapes:
            #         if not shape.has_text_frame:
            #             continue
            #         for paragraph in shape.text_frame.paragraphs:
            #             for run in paragraph.runs:
            #                 ret_value.append(run.text)

            # for txt in ret_value:
            #     print(txt)
            slide = self.presentation.slides[1]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    print(para.text)
                    for run in para.runs:
                        ret_value.append(run.text)
            # for txt in ret_value:
            #     print(txt)
            print("--------------x-------------------")
        except Exception as e:
            print("exception")
            self.exception_handler.handle_exception(e)

    def mine_text(self, array):

        pass

    def find_tables(self):
        slide = self.presentation.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                self.extract_text_from_table(shape.table)
        try:
            pass
        except Exception as e:
            print(e)

    def extract_text_from_table(self, table):
        table_data = {}
        try:
            previous_key = None
            for cell in table.iter_cells():
                if not previous_key:
                    previous_key = cell.text
                else:
                    table_data[previous_key] = cell.text
                    previous_key = None
            print(table_data)
        except Exception as e:
            print(e)


#minner = TextMinner("Recommended flights round 2 v1.pptx")
list_of_presentations = os.listdir("Requirements")
# print(list_of_presentations)


# for file_name in list_of_presentations:
#     minner = TextMinner(os.path.join("Requirements",file_name))
#     minner.extract_text()

minner = TextMinner("Requirements/Test.pptx")
table = minner.find_tables()
